from docx import Document
from pptx import Presentation
import os

#Import note file of various types
def openFile(file):
	filename, file_extension = os.path.splitext(file)
	if(file_extension==".docx"):
		document = Document(file)
		for para in document.paragraphs:
		    return para.text
	elif(file_extension==".pptx"):
		ppt = Presentation(file)
		for slide in ppt.slides:
			for shape in slide.shapes:
				print(shape.text)
	else:
		print("Cannot read filetype " + file_extension)

openFile("HELLO.pptx")
#Read it for Question templates
#Add it to dataFile.json